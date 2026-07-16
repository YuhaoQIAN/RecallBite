"""Parse uploaded documents into structured material."""

from __future__ import annotations

from io import BytesIO

from src.parsers.text_parser import MaterialDocument


SUPPORTED_EXTS = {".txt", ".md", ".pdf", ".docx", ".pptx"}
MAX_FILE_SIZE_MB = 10


def parse_document(file_bytes: bytes, filename: str) -> MaterialDocument:
    """Parse an uploaded document based on its extension."""
    if not file_bytes:
        raise ValueError("Empty file")

    size_mb = len(file_bytes) / (1024 * 1024)
    if size_mb > MAX_FILE_SIZE_MB:
        raise ValueError(f"File too large: {size_mb:.1f} MB (max {MAX_FILE_SIZE_MB} MB)")

    ext = filename.lower().split(".")[-1] if "." in filename else ""
    if ext == "txt" or ext == "md":
        text = file_bytes.decode("utf-8", errors="replace")
        return MaterialDocument(
            text=text.strip(),
            source_kind="uploaded_file",
            source_title=filename,
            source_reference=filename,
        )

    if ext == "pdf":
        return _parse_pdf(file_bytes, filename)

    if ext == "docx":
        return _parse_docx(file_bytes, filename)

    if ext == "pptx":
        return _parse_pptx(file_bytes, filename)

    raise ValueError(f"Unsupported file type: .{ext}")


def _parse_pdf(file_bytes: bytes, filename: str) -> MaterialDocument:
    """Parse PDF with PyMuPDF for coordinate-level extraction.
    
    Falls back to pypdf if PyMuPDF is not available.
    """
    # Try advanced PyMuPDF parser first
    try:
        from src.parsers.pdf_parser import parse_pdf_advanced
        parsed = parse_pdf_advanced(file_bytes, filename)
        
        if not parsed.full_text.strip():
            raise ValueError("PDF appears to be empty or contains no extractable text")
        
        # Detect overall language
        if parsed.is_bilingual:
            lang = "bilingual"
        elif "zh" in parsed.detected_languages:
            lang = "zh"
        elif "en" in parsed.detected_languages:
            lang = "en"
        else:
            lang = ""
        
        # Build structured pages for downstream use
        structured = []
        for sec in parsed.pages:
            structured.append({
                "page": sec.page,
                "language": sec.language,
                "title": sec.title,
                "text": sec.text,
                "column": sec.column,
            })
        
        return MaterialDocument(
            text=parsed.full_text.strip(),
            source_kind="uploaded_file",
            source_title=filename,
            source_reference=filename,
            location_info=f"{parsed.page_count} pages",
            structured_pages=structured,
            detected_language=lang,
            is_bilingual=parsed.is_bilingual,
        )
    except ImportError:
        pass  # Fall through to pypdf
    except Exception:
        pass  # Fall through to pypdf on any error
    
    # Fallback: pypdf
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ImportError("PDF parsing requires 'pypdf' or 'PyMuPDF'. Install: pip install pypdf PyMuPDF") from exc

    from io import BytesIO
    reader = PdfReader(BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text:
            pages.append(f"[Page {i + 1}]\n{page_text.strip()}")

    text = "\n\n".join(pages)
    if not text.strip():
        raise ValueError("PDF appears to be empty or contains no extractable text")

    return MaterialDocument(
        text=text.strip(),
        source_kind="uploaded_file",
        source_title=filename,
        source_reference=filename,
        location_info=f"{len(reader.pages)} pages",
    )


def _parse_docx(file_bytes: bytes, filename: str) -> MaterialDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise ImportError("DOCX parsing requires 'python-docx'. Install: pip install python-docx") from exc

    doc = Document(BytesIO(file_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    if not text.strip():
        raise ValueError("DOCX appears to be empty")

    return MaterialDocument(
        text=text.strip(),
        source_kind="uploaded_file",
        source_title=filename,
        source_reference=filename,
    )


def _parse_pptx(file_bytes: bytes, filename: str) -> MaterialDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("PPTX parsing requires 'python-pptx'. Install: pip install python-pptx") from exc

    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

    text = "\n\n".join(slides)
    if not text.strip():
        raise ValueError("PPTX appears to be empty or contains no text")

    return MaterialDocument(
        text=text.strip(),
        source_kind="uploaded_file",
        source_title=filename,
        source_reference=filename,
        location_info=f"{len(prs.slides)} slides",
    )
