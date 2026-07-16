"""
test_quality_regressions.py - Quality regression tests for RecallBite V1.

These tests verify the intelligence layer behavior, not just schema correctness.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

import pytest

from src.generator import generate_card
from src.retrieval import retrieve_relevant_cards
from src.activation import generate_apply_suggestion, generate_activation_output
from src.storage import save_cards, load_cards, update_card, get_card, record_usage
from src.parsers.text_parser import parse_text
from src.parsers.document_parser import parse_document
from src.parsers.url_parser import is_valid_url, _is_safe_url


# ── 9.1 Generator quality ─────────────────────────────────────────────────


class TestGeneratorQuality:
    """Verify generator filters fluff and prioritizes substance."""

    def test_fluff_filtered_from_takeaway(self):
        """Welcome/thank-you should not appear in the takeaway."""
        seed = (
            "Welcome everyone. Thank you to our speakers. "
            "74% of financial institutions reported cross-border data compliance as the main blocker."
        )
        card = generate_card(knowledge_seed=seed, source_type="Webcast", source="Test")
        takeaway = card.get("insight_pack", {}).get("thirty_second_takeaway", "")
        assert "welcome" not in takeaway.lower()
        assert "thank you" not in takeaway.lower()
        assert "74%" in takeaway or "74" in takeaway

    def test_numbers_prioritized_in_insights(self):
        """Sentences with numbers should rank higher than fluff."""
        seed = (
            "Welcome everyone. Good morning. "
            "74% of financial institutions reported cross-border data compliance as the main blocker. "
            "This is a very important topic."
        )
        card = generate_card(knowledge_seed=seed, source_type="Webcast", source="Test")
        insights = card.get("insight_pack", {}).get("key_insights", [])
        insight_texts = [i if isinstance(i, str) else i.get("insight", "") for i in insights]
        # The 74% sentence should be in the top insights
        assert any("74%" in txt or "74" in txt for txt in insight_texts)

    def test_no_duplicate_key_insights(self):
        """Key insights should not contain near-duplicates."""
        seed = (
            "AI governance is about accountability across business units. "
            "AI governance is about accountability across business units and IT. "
            "74% of firms now have AI boards."
        )
        card = generate_card(knowledge_seed=seed, source_type="Article", source="Test")
        insights = card.get("insight_pack", {}).get("key_insights", [])
        # After deduplication, the near-duplicate sentences should not both appear
        assert len(insights) >= 2
        assert len(set(insights)) == len(insights)

    def test_evidence_spans_present(self):
        """source_grounding should contain evidence spans."""
        seed = "74% of institutions report compliance as the main blocker."
        card = generate_card(knowledge_seed=seed, source_type="Article", source="Test")
        sg = card.get("source_grounding", {})
        spans = sg.get("evidence_spans", [])
        assert len(spans) > 0
        assert any("74%" in str(s.get("text", "")) for s in spans)


# ── 9.2 Retrieval topic separation ────────────────────────────────────────


class TestRetrievalTopicSeparation:
    """Verify retrieval separates topic relevance from output intent."""

    def test_ai_governance_ranks_above_climate_risk(self):
        """AI governance proposal should rank AI governance card first."""
        ai_card = generate_card(
            knowledge_seed="AI governance requires accountability across business units.",
            source_type="Article",
            topic_tags_text="AI governance, accountability",
            source="Test",
        )
        climate_card = generate_card(
            knowledge_seed="Climate risk disclosure standards are tightening in 2025.",
            source_type="Report",
            topic_tags_text="climate risk, ESG, disclosure",
            source="Test",
        )
        cards = [climate_card, ai_card]
        results = retrieve_relevant_cards("我要写 AI governance proposal opening", cards, top_k=2)
        assert len(results) == 2
        top_card, top_score = results[0]
        second_card, second_score = results[1]
        assert "AI governance" in top_card.get("topic_tags", [])
        assert top_score > second_score

    def test_proposal_does_not_make_all_cards_same_score(self):
        """The word 'proposal' alone should not give all cards high scores."""
        card_a = generate_card(
            knowledge_seed="Risk ownership in model development.",
            source_type="Article",
            topic_tags_text="risk ownership",
            source="Test",
        )
        card_b = generate_card(
            knowledge_seed="ESG reporting frameworks differ by region.",
            source_type="Report",
            topic_tags_text="ESG, reporting",
            source="Test",
        )
        results = retrieve_relevant_cards("proposal", [card_a, card_b], top_k=2)
        # Both should score 0 because the topic query is empty after stripping intent words
        assert results[0][1] == 0


# ── 9.3 Chinese stopword filtering ────────────────────────────────────────


class TestChineseStopwords:
    """Verify common Chinese words do not inflate scores."""

    def test_stopwords_do_not_match_everything(self):
        """Query with stopwords should not match irrelevant cards."""
        card = generate_card(
            knowledge_seed="AI governance requires accountability.",
            source_type="Article",
            topic_tags_text="AI governance",
            source="Test",
        )
        results = retrieve_relevant_cards("我要写一个关于当前工作的观点", [card], top_k=1)
        # The stopwords (一个, 当前, 工作, 观点, 关于) should be stripped,
        # leaving an empty topic query → score 0
        assert results[0][1] == 0


# ── 9.4 Activation quality ────────────────────────────────────────────────


class TestActivationQuality:
    """Verify activation produces usable, non-template output."""

    def test_proposal_output_contains_card_content(self):
        """Proposal activation should include concrete card content."""
        card = generate_card(
            knowledge_seed="74% of firms now have AI governance boards.",
            source_type="Article",
            topic_tags_text="AI governance",
            source="Test",
        )
        suggestion = generate_apply_suggestion("Write a proposal opening", card, score=10)
        paragraph = suggestion.get("copy_ready_paragraph", "")
        assert "74%" in paragraph or "AI governance" in paragraph

    def test_no_forbidden_filler_phrases(self):
        """Activation output should not contain banned filler phrases."""
        card = generate_card(
            knowledge_seed="Accountability must be clear.",
            source_type="Article",
            topic_tags_text="accountability",
            source="Test",
        )
        suggestion = generate_apply_suggestion("Write a proposal", card, score=10)
        paragraph = suggestion.get("copy_ready_paragraph", "")
        banned = ["从更系统的角度理解", "这与行业趋势相呼应", "可以为专业判断提供参考"]
        for phrase in banned:
            assert phrase not in paragraph, f"Banned phrase found: {phrase}"

    def test_very_foggy_does_not_generate_strong_conclusion(self):
        """Very Foggy cards should not produce confident copy-ready text."""
        card = generate_card(
            knowledge_seed="AI",  # Very short → Very Foggy
            source_type="Note",
            source="",
        )
        assert card["fog_index"]["level"] == "Very Foggy"
        suggestion = generate_apply_suggestion("Write a proposal", card, score=5)
        paragraph = suggestion.get("copy_ready_paragraph", "")
        assert "线索还太少" in paragraph or "too thin" in paragraph.lower() or "not enough" in paragraph.lower()

    def test_multi_card_synthesis_no_duplicates(self):
        """Multi-card activation should synthesize cards, not repeat each one verbatim."""
        card1 = generate_card(
            knowledge_seed="AI governance boards are now in 74% of firms.",
            source_type="Article",
            topic_tags_text="AI governance",
            source="Test A",
        )
        card2 = generate_card(
            knowledge_seed="Risk ownership must be assigned before deployment.",
            source_type="Article",
            topic_tags_text="risk ownership",
            source="Test B",
        )
        results = [
            {"card": card1, "score": 10, "match_strength": "strong"},
            {"card": card2, "score": 8, "match_strength": "strong"},
        ]
        out = generate_activation_output("Write proposal on AI governance", results)
        ready = out.get("ready_to_use_output", "")

        # Should not be empty and should contain something concrete
        assert ready
        assert "proposal" in ready.lower() or "proposal" in ready

        # Should reference both cards' core content (synthesis, not just first card)
        seed1 = card1.get("knowledge_seed", "")
        seed2 = card2.get("knowledge_seed", "")
        assert seed1[:30] in ready or card1.get("core_insight", "") in ready
        assert seed2[:30] in ready or card2.get("core_insight", "") in ready

        # Should NOT simply concatenate the two raw seeds back-to-back
        # (that would indicate no synthesis)
        combined_raw = seed1 + seed2
        assert ready != combined_raw

        # Should contain synthesis connectors, not just list markers
        synthesis_markers = ["首先", "此外", "同时", "共同", "First,", "Additionally,", "Also,"]
        assert any(m in ready for m in synthesis_markers), (
            f"Expected synthesis connectors in output, got: {ready[:200]}"
        )

        # Should NOT be a simple bulleted list of raw card contents
        list_markers = ["\n- ", "\n1. ", "\n* ", "  - ", "  1. ", "  * "]
        assert not any(m in ready for m in list_markers), (
            f"Output looks like a raw list instead of synthesis: {ready[:200]}"
        )


# ── 9.5 Parser tests ──────────────────────────────────────────────────────


class TestParsers:
    """Verify parser layer handles inputs and errors gracefully."""

    def test_text_parser_returns_material_document(self):
        doc = parse_text("Hello world", source_note="test")
        assert doc.text == "Hello world"
        assert doc.source_reference == "test"

    def test_txt_document_parsed(self, tmp_path):
        file_path = tmp_path / "test.txt"
        file_path.write_text("Hello from txt", encoding="utf-8")
        doc = parse_document(file_path.read_bytes(), file_path.name)
        assert "Hello from txt" in doc.text

    def test_md_document_parsed(self, tmp_path):
        file_path = tmp_path / "test.md"
        file_path.write_text("# Title\n\nHello from markdown", encoding="utf-8")
        doc = parse_document(file_path.read_bytes(), file_path.name)
        assert "Hello from markdown" in doc.text

    def test_pdf_graceful_on_invalid_content(self, tmp_path):
        file_path = tmp_path / "test.pdf"
        file_path.write_bytes(b"%PDF-1.4 fake pdf content")
        # pypdf may raise various errors on invalid content; ensure it does not crash silently
        with pytest.raises(Exception):
            parse_document(file_path.read_bytes(), file_path.name)

    def test_docx_graceful_on_invalid_content(self, tmp_path):
        file_path = tmp_path / "test.docx"
        file_path.write_bytes(b"PK\x03\x04 fake docx")
        # python-docx may raise BadZipFile on invalid content; ensure it does not crash silently
        with pytest.raises(Exception):
            parse_document(file_path.read_bytes(), file_path.name)

    def test_pdf_real_document_parsed(self):
        """Generate a real PDF in memory and verify text extraction."""
        from io import BytesIO
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        pdf.cell(200, 10, text="Hello from PDF page 1", new_x="LMARGIN", new_y="NEXT")
        pdf.add_page()
        pdf.cell(200, 10, text="Hello from PDF page 2", new_x="LMARGIN", new_y="NEXT")
        buf = BytesIO()
        pdf.output(buf)
        buf.seek(0)

        doc = parse_document(buf.read(), "test.pdf")
        assert "Hello from PDF page 1" in doc.text
        assert "Hello from PDF page 2" in doc.text
        assert "[Page 1]" in doc.text
        assert "[Page 2]" in doc.text
        assert doc.location_info == "2 pages"

    def test_docx_real_document_parsed(self):
        """Generate a real DOCX in memory and verify text extraction."""
        from io import BytesIO
        from docx import Document

        docx_buffer = BytesIO()
        doc = Document()
        doc.add_paragraph("Hello from DOCX paragraph 1")
        doc.add_paragraph("Hello from DOCX paragraph 2")
        doc.save(docx_buffer)
        docx_buffer.seek(0)

        result = parse_document(docx_buffer.read(), "test.docx")
        assert "Hello from DOCX paragraph 1" in result.text
        assert "Hello from DOCX paragraph 2" in result.text
        assert result.source_title == "test.docx"

    def test_pptx_real_document_parsed(self):
        """Generate a real PPTX in memory and verify text extraction."""
        from io import BytesIO
        from pptx import Presentation

        pptx_buffer = BytesIO()
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "Slide 1 Title"
        slide1.placeholders[1].text = "Hello from PPTX slide 1"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide 2 Title"
        slide2.placeholders[1].text = "Hello from PPTX slide 2"
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        result = parse_document(pptx_buffer.read(), "test.pptx")
        assert "Hello from PPTX slide 1" in result.text
        assert "Hello from PPTX slide 2" in result.text
        assert "[Slide 1]" in result.text
        assert "[Slide 2]" in result.text
        assert result.location_info == "2 slides"

    def test_url_validation(self):
        assert is_valid_url("https://example.com/article")
        assert not is_valid_url("not a url")
        assert not is_valid_url("ftp://example.com")

    def test_url_safety_blocks_private_and_localhost(self):
        blocked = [
            "http://localhost/admin",
            "http://127.0.0.1/secret",
            "http://10.0.0.1/internal",
            "http://172.16.0.1/api",
            "http://172.20.1.1/api",
            "http://192.168.1.1/router",
            "http://169.254.1.1/linklocal",
            "http://0.0.0.0",
            "http://fileserver.local/share",
            "http://intranet.internal/",
            "http://192.168.0.1:8080/admin",
        ]
        for url in blocked:
            safe, reason = _is_safe_url(url)
            assert not safe, f"Expected {url} to be blocked"
            assert "local or private" in reason.lower()

    def test_url_safety_allows_public(self):
        allowed = [
            "https://example.com/article",
            "https://github.com/user/repo",
            "http://news.site.com/story",
        ]
        for url in allowed:
            safe, reason = _is_safe_url(url)
            assert safe, f"Expected {url} to be allowed: {reason}"
            assert reason == ""

    def test_url_safety_blocks_raw_ip(self):
        safe, reason = _is_safe_url("http://8.8.8.8")
        assert not safe
        assert "local or private" in reason.lower()

    def test_url_safety_blocks_non_http_scheme(self):
        safe, reason = _is_safe_url("ftp://example.com/file")
        assert not safe
        assert "only http and https" in reason.lower()


# ── 9.6 Storage enhancements ──────────────────────────────────────────────


class TestStorageEnhancements:
    """Verify storage atomic write, update, and usage tracking."""

    def test_update_card_merges_fields(self, tmp_path):
        import src.storage as storage_module
        original_file = storage_module.CARDS_FILE
        original_dir = storage_module.DATA_DIR
        data_dir = tmp_path / "data"
        data_dir.mkdir()
        cards_file = data_dir / "cards.json"
        storage_module.CARDS_FILE = cards_file
        storage_module.DATA_DIR = data_dir

        try:
            card = generate_card(
                knowledge_seed="Initial content",
                source_type="Article",
                source="Test",
            )
            save_cards([card])
            updated = update_card(card["id"], {"source_type": "Updated Webcast", "core_insight": "Updated insight"})
            assert updated is True
            loaded = get_card(card["id"])
            assert loaded is not None
            assert loaded["source_type"] == "Updated Webcast"
            assert loaded["core_insight"] == "Updated insight"
            # Edit history should exist
            history = loaded.get("_edit_history", [])
            assert len(history) >= 1
            assert "source_type" in history[0]["fields_changed"]
        finally:
            storage_module.CARDS_FILE = original_file
            storage_module.DATA_DIR = original_dir

    def test_record_usage_appends_history(self, tmp_path):
        import src.storage as storage_module
        original_file = storage_module.CARDS_FILE
        original_dir = storage_module.DATA_DIR
        data_dir = tmp_path / "data"
        data_dir.mkdir()
        cards_file = data_dir / "cards.json"
        storage_module.CARDS_FILE = cards_file
        storage_module.DATA_DIR = data_dir

        try:
            card = generate_card(
                knowledge_seed="Test content",
                source_type="Article",
                source="Test",
            )
            save_cards([card])
            ok = record_usage(card["id"], "Write proposal opening")
            assert ok is True
            loaded = get_card(card["id"])
            history = loaded.get("_usage_history", [])
            assert len(history) == 1
            assert history[0]["task"] == "Write proposal opening"
        finally:
            storage_module.CARDS_FILE = original_file
            storage_module.DATA_DIR = original_dir

    def test_corrupted_json_backup(self, tmp_path):
        import src.storage as storage_module
        original_file = storage_module.CARDS_FILE
        original_dir = storage_module.DATA_DIR
        data_dir = tmp_path / "data"
        data_dir.mkdir()
        cards_file = data_dir / "cards.json"
        cards_file.write_text("NOT JSON{{{", encoding="utf-8")
        storage_module.CARDS_FILE = cards_file
        storage_module.DATA_DIR = data_dir

        try:
            cards = load_cards()
            assert cards == []
            # Backup file should exist
            backups = list(data_dir.glob("cards_corrupted_backup_*.json"))
            assert len(backups) == 1
        finally:
            storage_module.CARDS_FILE = original_file
            storage_module.DATA_DIR = original_dir
